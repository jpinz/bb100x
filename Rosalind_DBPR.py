from Bio import ExPASy
from Bio import SwissProt
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE

# Get info about a protein from uniprot.org

protein_id = input("What is the UniProt ID?: ")
handle = ExPASy.get_sprot_raw(protein_id)
record = SwissProt.read(handle)

protein = record.description.split("=")[1].split("{")[0].strip()
name = record.organism.split("(")[0].strip()
fcn = record.comments[0].split(":")[1].split("{")[0].strip()

# Get the functions of the protein
processes = []
for x in record.cross_references:
    if x[0] == 'GO':
        if x[2][0] == 'P':
            processes.append(x[2][2:])

# Create a ppt
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

# Powerpoint slide formatting
title_shape.text = '{}'.format(name)
title_shape.word_wrap = False
title_shape.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

prot = body_shape.text_frame
prot.text = 'Protein: {}'.format(protein)
desc = prot.add_paragraph()
desc.text = "Description: {}".format(fcn)

fcns = prot.add_paragraph()
fcns.text = 'Functions:'

# add the functions to the list
for x in processes:
    item = prot.add_paragraph()
    item.text = x
    item.level = 1

# Save ppt
prs.save('{}.pptx'.format(protein.strip()))
